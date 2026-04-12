"""
Create a PPTX file with the Banking Governance & Quality Flow diagram + explanatory slides
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 25, 112)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(173, 216, 230)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_points):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(25, 25, 112)
    
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(9), Inches(0))
    line.line.color.rgb = RGBColor(70, 130, 180)
    line.line.width = Pt(3)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_before = Pt(6)
        p.space_after = Pt(6)

def add_diagram_description_slide(prs, title, description_text):
    """Add a slide describing the flow diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(25, 25, 112)
    
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(9), Inches(0))
    line.line.color.rgb = RGBColor(70, 130, 180)
    line.line.width = Pt(3)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = description_text
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.line_spacing = 1.4

# Slide 1: Title
add_title_slide(prs, "Banking Data Flow", "Governance & Quality - End-to-End")

# Slide 2: Overview of the Flow
add_diagram_description_slide(prs, "Data Flow: Overview",
"""The diagram shows 8 interconnected layers:

1. SOURCES (Blue) – Core banking, transactions, CRM, loans
2. GOVERNANCE (Purple) – Policies (PCI, AML, GDPR), SLAs, data owners
3. ETL/INGEST (Green) – Normalize, deduplicate, store
4. DATA WAREHOUSE (Pink) – Customers, Accounts, Transactions, Loans tables
5. QUALITY CHECKS (Orange) – Validation rules, monitoring, alerts
6. REMEDIATION (Light Green) – Auto-fixes, manual review, incident tracking
7. COMPLIANCE (Light Purple) – AML checks, PII protection, audit trails
8. CONSUMPTION (Teal) – Reports, dashboards, fraud detection, segmentation
9. FEEDBACK LOOP (Yellow) – Issues, SLA updates, continuous improvement
""")

# Slide 3: Banking Sources
add_content_slide(prs, "Step 1: Data Sources (Blue Layer)", [
    "💳 Core Banking System – Accounts, balances, account types",
    "📊 Transaction Log – Deposits, withdrawals, transfers",
    "👤 CRM System – Customer KYC, profile, identity info",
    "📋 Loan Management – Applications, repayments, status",
    "⚠️ Challenge: Data from disparate legacy systems, different formats/frequencies"
])

# Slide 4: Governance Layer
add_content_slide(prs, "Step 2: Governance Policies (Purple Layer)", [
    "🔐 PCI-DSS – Card data encrypted; access audited",
    "🚨 AML/KYC – Sanction screening; transaction monitoring",
    "📋 GDPR – Customer consent; data retention (7 years for txns)",
    "🎯 Data Owners – Finance (accounts), Compliance (KYC), Risk (loans)",
    "⏱️ SLAs – Account balance <5 min latency; fraud alerts <10 sec; KYC <1 day"
])

# Slide 5: ETL/Ingest
add_content_slide(prs, "Step 3: ETL & Ingestion (Green Layer)", [
    "1️⃣ Ingest & Normalize – Standardize formats from disparate systems",
    "2️⃣ Deduplication – Match customers across systems (SSN + DOB)",
    "3️⃣ PII Masking – Early masking of sensitive data (card numbers, SSN)",
    "4️⃣ Store – Load into warehouse (Snowflake, BigQuery) with partitioning",
    "✅ Outcome: Single source of truth; audit trail preserved"
])

# Slide 6: Data Warehouse Entities
add_content_slide(prs, "Step 4: Data Warehouse Entities (Pink Layer)", [
    "👥 Customers – ID, Name, KYC_Status, Risk_Level",
    "💰 Accounts – Account_ID, Customer_ID, Balance, Type, Last_Updated",
    "📊 Transactions – Txn_ID, Account_ID, Amount, Date, Type, Status",
    "📋 Loans – Loan_ID, Customer_ID, Amount, Status, Start_Date",
    "🔗 Relationships: Foreign keys ensure referential integrity; lineage tracked"
])

# Slide 7: Quality Rules
add_content_slide(prs, "Step 5: Data Quality Rules (Orange Layer)", [
    "✓ Customers – No null KYC; unique SSN; valid email",
    "✓ Accounts – Balance >= 0; matched to customer; no orphans",
    "✓ Transactions – Amount > 0; date <= today; reconciles to GL; account exists",
    "✓ Loans – Amount > 0; customer verified; start_date valid",
    "📈 Monitoring – Real-time: transaction count trends, reconciliation %, duplicates flagged"
])

# Slide 8: Remediation
add_content_slide(prs, "Step 6: Data Remediation (Light Green Layer)", [
    "🤖 Automated Fixes – Trim/standardize names, dedup by SSN+DOB, backfill dates",
    "👀 Manual Review – Conflicting records, suspicious transactions, high-risk customers",
    "🎫 Incident Tracking – Create ticket, track MTTR, RCA, prevention rules",
    "🔄 Reprocess – Backfill corrected data to meet SLAs",
    "📋 Outcome: Continuous quality improvement; incidents resolved; SLAs met"
])

# Slide 9: Compliance & Audit
add_content_slide(prs, "Step 7: Compliance & Audit (Light Purple Layer)", [
    "🚨 AML – Sanction screening, transaction monitoring, SAR generation",
    "🔐 PII Protection – Mask SSN/account# for analysts; encrypt at rest",
    "📋 Audit Trail – Who accessed what, when, for what purpose; data lineage",
    "📊 Reporting – Regulatory compliance checked; evidence collected",
    "✅ Outcome: Regulatory mandates met; risks mitigated; compliance demonstrated"
])

# Slide 10: Consumption
add_content_slide(prs, "Step 8: Consumption & Applications (Teal Layer)", [
    "📄 Financial Reports – Balance sheets, P&L for regulators & execs",
    "⚠️ Risk Dashboard – Credit exposure, NPL %, capital ratios",
    "🚨 Fraud Detection – ML anomaly models detect suspicious patterns <10 sec",
    "🎯 Customer Segmentation – Lifetime value, churn risk for marketing",
    "✅ Outcome: Insights for decisioning; fast fraud prevention; strategic alignment"
])

# Slide 11: Feedback Loop
add_content_slide(prs, "Step 9: Feedback & Continuous Improvement (Yellow Layer)", [
    "❌ Report Issues – 'Balance wrong in report', 'missed fraud', 'late transactions'",
    "🔧 Update SLAs – Tighten thresholds, add new quality rules based on incidents",
    "🔄 Loop Back to Governance – New policies, owner responsibilities cascade down",
    "📊 Measure Impact – Track MTTR, incident count, DQ score improvements",
    "✅ Outcome: Continuous learning; faster issue resolution; proactive governance"
])

# Slide 12: Key Players & Responsibilities
add_content_slide(prs, "Key Roles & Responsibilities", [
    "👔 Data Owner – Defines policies, SLAs, assigns stewards (e.g., Finance owns Accounts)",
    "⚙️ Data Engineer – Builds pipelines, ETL, monitoring, infrastructure",
    "📋 Data Steward – Monitors quality, responds to alerts, roots cause issues",
    "🔒 Compliance Officer – Ensures regulatory controls, audits, reports",
    "👨‍💼 Analyst / Risk Officer – Consumes data, flags issues, provides feedback"
])

# Slide 13: Success Metrics
add_content_slide(prs, "KPIs to Track Success", [
    "📈 DQ Score (0–100) – Target: >95% across critical datasets",
    "% Failing Records – Target: <1% in production",
    "⏱️ MTTR (Mean Time to Resolve) – Target: <4 hours for data incidents",
    "✅ SLA Compliance – Target: >99% of datasets meet freshness SLA",
    "🎯 Data Catalog Coverage – Target: 100% of critical datasets documented"
])

# Slide 14: Common Pitfalls
add_content_slide(prs, "Common Pitfalls & How to Avoid", [
    "❌ No clear ownership → Assign domain stewards with clear escalation",
    "❌ Over-automation without checks → Add human review for ambiguous fixes",
    "❌ One-time project mindset → Make it ongoing program with sponsorship",
    "❌ Poor metadata → Prioritize catalog + minimal required fields",
    "❌ Siloed teams → Cross-functional governance council + regular syncs"
])

# Slide 15: Implementation Timeline
add_content_slide(prs, "Quick Implementation (90 Days)", [
    "📅 Week 1–2: Inventory top 5 datasets; run profiling; understand baseline",
    "📅 Week 3–4: Assign owners per domain; define SLAs & data contracts",
    "📅 Week 5–8: Implement validation tests (dbt, Great Expectations); add to pipelines",
    "📅 Week 9–10: Build dashboards; configure alerts for SLA violations",
    "📅 Week 11–12: Scale to more datasets; formalize governance council; training"
])

# Slide 16: Conclusion
add_content_slide(prs, "Key Takeaways", [
    "✅ Data governance enforces policies; quality rules ensure data usability",
    "✅ In banking, speed (fraud <10 sec) and compliance (PCI, AML) are critical",
    "✅ Every data issue cascades downstream – lineage and feedback loops matter",
    "✅ Success requires cross-functional ownership and continuous improvement",
    "✅ Tools enable automation, but people and processes drive results"
])

# Save
output_path = r"e:\VScode26\Banking_Data_Flow_Governance_Quality.pptx"
prs.save(output_path)
print(f"✅ PPTX created: {output_path}")
print(f"   16 slides covering the complete banking data governance & quality flow")
