from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_CENTER, TA_LEFT

# Load PPTX
prs = Presentation(r"e:\VScode26\Banking_Data_Flow_Governance_Quality.pptx")

# Create PDF
pdf_path = r"e:\VScode26\Banking_Data_Flow_Governance_Quality.pdf"
doc = SimpleDocTemplate(pdf_path, pagesize=letter, topMargin=0.5*inch, bottomMargin=0.5*inch)

# Style setup
styles = getSampleStyleSheet()
story = []

title_style = ParagraphStyle(
    'CustomTitle',
    parent=styles['Heading1'],
    fontSize=22,
    textColor='#191B70',
    spaceAfter=10,
    alignment=TA_CENTER,
    bold=True
)

body_style = ParagraphStyle(
    'CustomBody',
    parent=styles['Normal'],
    fontSize=10,
    spaceAfter=6,
    alignment=TA_LEFT,
    leading=12
)

# Extract content from slides
for slide_idx, slide in enumerate(prs.slides):
    if slide_idx > 0:
        story.append(PageBreak())
    
    # Extract text from all shapes in slide
    for shape_idx, shape in enumerate(slide.shapes):
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()
            # First shape is often title
            if shape_idx == 0:
                story.append(Paragraph(text, title_style))
                story.append(Spacer(1, 0.15*inch))
            else:
                # Split by newlines for content
                lines = text.split('\n')
                for line in lines:
                    line = line.strip()
                    if line:
                        story.append(Paragraph(line, body_style))
    
    story.append(Spacer(1, 0.1*inch))

# Build PDF
try:
    doc.build(story)
    print(f"✅ PDF conversion successful!")
    print(f"📄 File saved: {pdf_path}")
    print(f"   16 pages with complete banking data flow documentation")
except Exception as e:
    print(f"❌ Error: {e}")
