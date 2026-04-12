"""
Create PPTX with Banking Data Flow Diagram (from Mermaid flowchart) 
with visual representation and detailed slides
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 25, 112)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(173, 216, 230)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_points):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(25, 25, 112)
    
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(9), Inches(0))
    line.line.color.rgb = RGBColor(70, 130, 180)
    line.line.width = Pt(3)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_before = Pt(6)
        p.space_after = Pt(6)

def add_diagram_overview_slide(prs):
    """Add a slide with visual boxes representing the flow layers"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Banking Data Governance & Quality Flow - Complete Overview"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(25, 25, 112)
    
    # Define layers with colors
    layers = [
        ("1. SOURCES", "Banking Data\nSources", Inches(0.5), Inches(1.0), RGBColor(225, 245, 255)),
        ("2. GOVERNANCE", "Policies &\nCompliance", Inches(0.5), Inches(1.8), RGBColor(243, 229, 245)),
        ("3. ETL", "Ingest &\nNormalize", Inches(0.5), Inches(2.6), RGBColor(232, 245, 233)),
        ("4. WAREHOUSE", "Data Storage\nTables", Inches(0.5), Inches(3.4), RGBColor(252, 228, 236)),
        ("5. QUALITY", "Validation &\nMonitoring", Inches(0.5), Inches(4.2), RGBColor(255, 243, 224)),
        ("6. REMEDIATION", "Fix & Improve\nData", Inches(0.5), Inches(5.0), RGBColor(241, 248, 233)),
        ("7. COMPLIANCE", "AML & Security\nChecks", Inches(0.5), Inches(5.8), RGBColor(237, 231, 246)),
        ("8. CONSUMPTION", "Analytics &\nReports", Inches(0.5), Inches(6.6), RGBColor(224, 242, 241)),
    ]
    
    for label, text, left, top, color in layers:
        # Add rectangle
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.1), Inches(0.65))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(70, 130, 180)
        shape.line.width = Pt(2)
        
        # Add label text
        text_frame = shape.text_frame
        text_frame.text = label
        text_frame.paragraphs[0].font.size = Pt(10)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = 1  # Middle
        
        # Add description to the right
        desc_box = slide.shapes.add_textbox(left + Inches(1.2), top, Inches(8), Inches(0.65))
        desc_frame = desc_box.text_frame
        desc_frame.text = text
        desc_frame.paragraphs[0].font.size = Pt(11)
        desc_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        desc_frame.vertical_anchor = 1

# Slide 1: Title
add_title_slide(prs, "Complete Banking Data Flow", "Governance & Quality - Step-by-Step")

# Slide 2: Diagram Overview
add_diagram_overview_slide(prs)

# Slide 3: Layer 1 - Banking Sources
add_content_slide(prs, "Layer 1: Banking Data Sources (Blue)", [
    "💳 Core Banking System – Accounts, balances, account types, account holders",
    "📊 Transaction Log – Deposits, withdrawals, transfers, interest payments",
    "👤 CRM System – Customer KYC data, profile, identity verification",
    "📋 Loan Management – Applications, approval status, repayments, interest",
    "⚠️ Challenge: Legacy systems with different formats, update frequencies, data quality"
])

# Slide 4: Layer 2 - Governance
add_content_slide(prs, "Layer 2: Governance: Policies & Compliance (Purple)", [
    "🔐 PCI-DSS – Card data encrypted; restricted access; audit logs",
    "🚨 AML/KYC – Sanction screening; customer due diligence; transaction monitoring",
    "📋 GDPR – Customer consent; data retention (7 years for transactions); PII handling",
    "🎯 Data Owners – Finance (accounts), Compliance (KYC), Risk (loans)",
    "⏱️ SLAs – Account <5min latency; fraud alerts <10 sec; KYC updates <1 day"
])

# Slide 5: Layer 3 - ETL/Ingest
add_content_slide(prs, "Layer 3: ETL/Ingest Pipeline (Green)", [
    "1️⃣ Ingest & Normalize – Standardize date formats, encodings, units across systems",
    "2️⃣ PII Masking – Early masking of SSN, card numbers, account details",
    "3️⃣ Deduplication & Customer Matching – Match same customer across systems (SSN+DOB)",
    "4️⃣ Store in Warehouse – Load to Snowflake/BigQuery with proper partitioning",
    "✅ Outcome: Single source of truth; data lineage preserved; audit trail complete"
])

# Slide 6: Layer 4 - Data Warehouse
add_content_slide(prs, "Layer 4: Data Warehouse - Banking Entities (Pink)", [
    "👥 Customers Table – ID, Name, KYC_Status, Risk_Level, Last_Updated",
    "💰 Accounts Table – Account_ID, Customer_ID, Balance, Type, Status",
    "📊 Transactions Table – Txn_ID, Account_ID, Amount, Date, Type, Status",
    "📋 Loans Table – Loan_ID, Customer_ID, Amount, Start_Date, Status",
    "🔗 Relationships: Foreign keys maintained; lineage tracked; data documented"
])

# Slide 7: Layer 5 - Quality Checks
add_content_slide(prs, "Layer 5: Data Quality Rules & Monitoring (Orange)", [
    "✓ Customers – No null KYC; unique SSN; valid email; active status",
    "✓ Accounts – Balance >= 0; customer ID exists; no orphaned records",
    "✓ Transactions – Amount > 0; date <= today; reconciles to General Ledger",
    "✓ Loans – Amount > 0; customer verified; start_date valid",
    "📈 Monitoring – Real-time metrics: transaction trends, reconciliation %, duplicates"
])

# Slide 8: Layer 6 - Remediation
add_content_slide(prs, "Layer 6: Data Remediation (Light Green)", [
    "🤖 Automated Fixes – Trim/standardize names; dedup by SSN+DOB; backfill missing dates",
    "👀 Manual Review – Conflicting customer records; suspicious transactions; high-risk flags",
    "🎫 Incident Tracking – Create tickets; track MTTR; root cause analysis; prevention rules",
    "🔄 Reprocess & Backfill – Update historical data; verify fixes; validate SLAs met",
    "✅ Outcome: Data corrected; incidents reduced; SLAs consistently met"
])

# Slide 9: Layer 7 - Compliance & Audit
add_content_slide(prs, "Layer 7: Compliance & Audit (Light Purple)", [
    "🚨 AML Monitoring – Sanction list screening; transaction pattern analysis; SAR filing",
    "🔐 PII Protection – Mask SSN/account# in reports; encrypt at rest; access logs",
    "📋 Audit Trail – Who accessed what, when, why; data lineage end-to-end",
    "📊 Regulatory Reporting – Accurate balance sheets; NPL %; capital ratios to central bank",
    "✅ Outcome: Regulatory compliance; risks managed; evidence trail complete"
])

# Slide 10: Layer 8 - Consumption
add_content_slide(prs, "Layer 8: Analytics & Consumer Applications (Teal)", [
    "📄 Financial Reports – Balance sheets, P&L, cash flow for regulators and executives",
    "⚠️ Risk Dashboard – Credit exposure, NPL %, capital adequacy ratios, stress testing",
    "🚨 Fraud Detection – ML model detects anomalies; alerts within 10 seconds",
    "🎯 Customer Segmentation – Lifetime value, churn risk prediction, marketing targeting",
    "✅ Outcome: Insights drive decisions; risks identified fast; competitive advantage"
])

# Slide 11: Feedback Loop
add_content_slide(prs, "Feedback Loop: Continuous Improvement (Yellow)", [
    "❌ Report Issues – Analysts flag 'wrong balance', 'missed fraud', 'late transactions'",
    "🔧 Investigate & Remediate – Data steward identifies root cause; fixes data",
    "📝 Update Policies – Tighten validation thresholds; add new quality rules; adjust SLAs",
    "🔄 Loop Back to Governance – New rules cascade down; all layers adapt",
    "✅ Outcome: Continuous learning; proactive governance; faster issue resolution"
])

# Slide 12: Key Connections
add_content_slide(prs, "How Layers Connect: The Flow", [
    "Sources → Governance defines policies that must be enforced",
    "Governance → ETL implements early controls (masking, dedup, validation)",
    "ETL → Warehouse stores clean, trusted data with lineage",
    "Warehouse → Quality rules validate data continuously",
    "Quality → Remediation fixes issues; Compliance audits; Consumption uses data",
    "Feedback → Issues reported; rules updated; cycle repeats"
])

# Slide 13: Risk Points & Mitigations
add_content_slide(prs, "Critical Risk Points in Banking", [
    "🔴 Wrong balance in account – Impact: Customer confusion, loss of trust",
    "   → Mitigation: GL reconciliation; balance >= 0 rule; daily audits",
    "🔴 Missed fraud – Impact: Financial loss, regulatory fines",
    "   → Mitigation: AML monitoring; anomaly detection; <10 sec alerts",
    "🔴 Customer KYC not updated – Impact: AML breach, regulatory penalty",
    "   → Mitigation: KYC SLA <1 day; automated escalation; audit trail"
])

# Slide 14: Success Metrics
add_content_slide(prs, "KPIs: Measuring Success", [
    "📈 DQ Score – Target: >95% across all critical datasets",
    "% Failing Records – Target: <1% of production records",
    "⏱️ MTTR (Mean Time to Resolve) – Target: <4 hours for data incidents",
    "✅ SLA Compliance – Target: >99% of datasets meet freshness/latency SLA",
    "🎯 Incident Rate – Target: <5 incidents per month from data quality issues"
])

# Slide 15: Implementation Timeline
add_content_slide(prs, "Quick Win: 90-Day Implementation", [
    "Week 1-2: Assess – Profile top 5 datasets; capture baseline DQ metrics",
    "Week 3-4: Prioritize – Assign owners; define SLAs per domain",
    "Week 5-8: Automate – Implement validation tests; set up monitoring dashboards",
    "Week 9-10: Remediate – Fix known issues; establish incident workflows",
    "Week 11-12: Scale – Expand to more datasets; formalize governance council"
])

# Slide 16: Conclusion
add_content_slide(prs, "Key Takeaways", [
    "✅ Data governance enforces policies; data quality ensures usability",
    "✅ Speed (fraud <10 sec) and compliance (PCI, AML) are non-negotiable in banking",
    "✅ Every layer impacts downstream – lineage and feedback loops matter",
    "✅ Success requires cross-functional ownership: finance, compliance, risk, engineering",
    "✅ Use tools to automate, but people and processes drive lasting results"
])

# Save
output_path = r"e:\VScode26\Banking_Data_Flow_Diagram_Complete.pptx"
prs.save(output_path)
print(f"✅ PPTX created: {output_path}")
print(f"   16 slides with complete banking data flow diagram breakdown")
