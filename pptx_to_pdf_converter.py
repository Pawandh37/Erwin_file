from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_CENTER, TA_LEFT

# Load PPTX
prs = Presentation(r"e:\VScode26\Data_Governance_Quality_Banking.pptx")

# Create PDF
pdf_path = r"e:\VScode26\Banking_Governance.pdf"
doc = SimpleDocTemplate(pdf_path, pagesize=letter)

# Style setup
styles = getSampleStyleSheet()
story = []

title_style = ParagraphStyle(
    'CustomTitle',
    parent=styles['Heading1'],
    fontSize=24,
    textColor='#191B70',
    spaceAfter=12,
    alignment=TA_CENTER,
    bold=True
)

heading_style = ParagraphStyle(
    'CustomHeading',
    parent=styles['Heading2'],
    fontSize=18,
    textColor='#191B70',
    spaceAfter=8,
    bold=True
)

body_style = ParagraphStyle(
    'CustomBody',
    parent=styles['Normal'],
    fontSize=11,
    spaceAfter=8,
    alignment=TA_LEFT
)

# Extract content from slides
for slide_idx, slide in enumerate(prs.slides):
    if slide_idx > 0:
        story.append(PageBreak())
    
    # Extract text from all shapes in slide
    slide_content = []
    for shape_idx, shape in enumerate(slide.shapes):
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()
            # First shape is often title
            if shape_idx == 0:
                story.append(Paragraph(text, title_style))
                story.append(Spacer(1, 0.2*inch))
            else:
                # Split by newlines for bullet points
                lines = text.split('\n')
                for line in lines:
                    line = line.strip()
                    if line:
                        story.append(Paragraph(f"• {line}", body_style))
    
    if slide_content:
        story.append(Spacer(1, 0.1*inch))

# Build PDF
try:
    doc.build(story)
    print(f"✅ PDF conversion successful!")
    print(f"📄 File saved: {pdf_path}")
except Exception as e:
    print(f"❌ Error: {e}")
