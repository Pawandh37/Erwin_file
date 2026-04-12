from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 25, 112)  # Midnight blue
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(173, 216, 230)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_points):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(25, 25, 112)
    
    # Add horizontal line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(9), Inches(0))
    line.line.color.rgb = RGBColor(70, 130, 180)
    line.line.width = Pt(3)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_before = Pt(8)
        p.space_after = Pt(8)

def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points):
    """Add a two-column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(25, 25, 112)
    
    # Left column
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.2), Inches(0.5))
    left_title_frame = left_title_box.text_frame
    left_title_frame.text = left_title
    left_title_frame.paragraphs[0].font.size = Pt(20)
    left_title_frame.paragraphs[0].font.bold = True
    left_title_frame.paragraphs[0].font.color.rgb = RGBColor(70, 130, 180)
    
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.2), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, point in enumerate(left_points):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(16)
        p.space_before = Pt(6)
    
    # Right column
    right_title_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.3), Inches(4.2), Inches(0.5))
    right_title_frame = right_title_box.text_frame
    right_title_frame.text = right_title
    right_title_frame.paragraphs[0].font.size = Pt(20)
    right_title_frame.paragraphs[0].font.bold = True
    right_title_frame.paragraphs[0].font.color.rgb = RGBColor(70, 130, 180)
    
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.9), Inches(4.2), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, point in enumerate(right_points):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(16)
        p.space_before = Pt(6)

# Slide 1: Title
add_title_slide(prs, "Data Governance & Data Quality", "Banking Industry Focus")

# Slide 2: Introduction
add_content_slide(prs, "Why Data Governance & Quality?", [
    "✓ Bad data leads to wrong decisions, model drift, and compliance risk",
    "✓ Regulatory mandates: PCI-DSS, AML/KYC, GDPR require data controls",
    "✓ Operational cost: data incidents cause delays, customer impact, fines",
    "✓ Trustworthiness: clean, well-documented data drives analytics and ML success",
    "✓ Competitive edge: organizations with strong data governance move faster"
])

# Slide 3: Data Quality Dimensions
add_content_slide(prs, "Data Quality: 7 Key Dimensions", [
    "1. Accuracy – Values correctly reflect real-world entities",
    "2. Completeness – Required fields and records are present",
    "3. Consistency – Same data agrees across systems and formats",
    "4. Timeliness – Data is current enough for its intended use",
    "5. Validity – Values conform to allowed formats and ranges",
    "6. Uniqueness – No unintended duplicates in key fields",
    "7. Integrity – Foreign keys and lineage are preserved"
])

# Slide 4: Data Governance Pillars
add_two_column_slide(prs, "Data Governance: Two Sides of the Coin",
    "Policy & Compliance",
    ["Policies & standards", "Metadata & catalog", "Retention & classification", "Access controls", "Audit trails"],
    "Operational Execution",
    ["Data owners & stewards", "Quality rules & testing", "Monitoring & alerts", "Remediation workflows", "Incident tracking"]
)

# Slide 5: Data Quality Workflow
add_content_slide(prs, "Data Quality Workflow", [
    "1️⃣ Profiling & Baseline – Scan data to understand distributions, nulls, outliers",
    "2️⃣ Rules & Tests – Codify validation checks (type, range, referential integrity)",
    "3️⃣ Monitoring & Alerts – Stream metrics; alert on SLA violations or anomalies",
    "4️⃣ Remediation – Auto-fixes for predictable issues; manual review for complex cases",
    "5️⃣ Backfill & Reprocess – Update historical data; verify fixes before marking resolved",
    "6️⃣ Feedback – Consumers report issues; priorities filter back to governance"
])

# Slide 6: Banking Data Entities
add_content_slide(prs, "Banking Data: Key Entities", [
    "💳 Customers – KYC status, identity, risk level; must be complete and unique",
    "💰 Accounts – Balance, type, owner; balance >= 0, must match to customer",
    "📊 Transactions – Amount, date, type; must reconcile to General Ledger",
    "📋 Loans – Amount, status, terms; amount > 0, customer verified",
    "⚠️ Core challenge: Data from disparate legacy systems → dedup, match, reconcile"
])

# Slide 7: Banking Quality Rules
add_two_column_slide(prs, "Banking: Data Quality Rules",
    "Customers & Accounts",
    ["No null KYC status", "Valid email format", "Unique SSN + DOB", "Balance >= 0", "No orphaned accounts"],
    "Transactions & Loans",
    ["Amount > 0", "Date <= today", "Account must exist", "GL reconciliation", "Loan start date valid"]
)

# Slide 8: Banking Compliance
add_content_slide(prs, "Banking: Compliance & Risk", [
    "🔐 PCI-DSS – Card data encrypted; access audited and logged",
    "🚨 AML/KYC – Sanction screening; transaction monitoring; suspicious activity reports (SARs)",
    "📋 GDPR – Customer consent tracked; data retention enforced; PII masked and protected",
    "📊 Regulatory Reporting – Accurate balances, NPL %, capital ratios sent to central bank",
    "⏱️ SLAs – Account balance <5 min latency; fraud alerts <10 sec; KYC updates <1 day"
])

# Slide 9: Banking Data Flow
add_content_slide(prs, "Banking Data Flow: End-to-End", [
    "1. Sources: Core banking, transaction log, CRM, loan systems",
    "2. Ingest & normalize: Standardize formats; mask PII early",
    "3. Deduplicate: Match customers across systems (SSN + DOB)",
    "4. Warehouse: Load into Snowflake/BigQuery",
    "5. Quality checks: Validate balances, reconcile GL, flag duplicates",
    "6. Compliance: Monitor for AML, audit access",
    "7. Consumption: Reports, dashboards, fraud detection, segmentation"
])

# Slide 10: Implementation Roadmap
add_content_slide(prs, "Implementation Roadmap: 90 Days", [
    "Week 1-2: Inventory & Prioritize – List top 5 business-critical datasets; run profiling",
    "Week 3-4: Assign Owners – Data owner for each domain; define SLAs and contracts",
    "Week 5-8: Implement Rules – Write validation tests (dbt, Great Expectations); add to pipelines",
    "Week 9-10: Monitoring Setup – Build dashboards; configure alerts for SLA violations",
    "Week 11-12: Scale & Iterate – Expand to more datasets; formalize governance council"
])

# Slide 11: Tools & Platforms
add_two_column_slide(prs, "Tools & Platforms",
    "Data Quality",
    ["Great Expectations", "Deequ", "Soda", "dbt tests", "Monte Carlo"],
    "Governance & Catalog",
    ["Microsoft Purview", "Collibra", "Alation", "Amundsen", "Informatica"]
)

# Slide 12: KPIs & Success Metrics
add_content_slide(prs, "KPIs to Track Progress", [
    "📈 DQ Score – Target: >95% across all critical datasets",
    "📊 % Failing Rules – Target: <1% of records in production",
    "⏱️ MTTR (Mean Time to Resolve) – Target: <4 hours for data incidents",
    "✅ SLA Compliance – Target: >99% of datasets meet freshness/latency SLA",
    "🎯 Data Catalog Coverage – Target: 100% of critical datasets documented and tagged"
])

# Slide 13: Common Pitfalls & Mitigations
add_two_column_slide(prs, "Common Pitfalls",
    "Problem",
    ["No clear ownership", "Over-automation", "Single project mindset", "Poor metadata"],
    "Mitigation",
    ["Assign domain stewards", "Add human reviews", "Ongoing program", "Catalog first"]
)

# Slide 14: Key Takeaways
add_content_slide(prs, "Key Takeaways", [
    "✅ Data governance is a continuous program, not a one-time project",
    "✅ Data quality is enforced through policies, tests, and monitoring—not just ideals",
    "✅ In banking, speed (fraud <10 sec) and compliance (AML, PCI) are non-negotiable",
    "✅ Every data issue has downstream impact—lineage and feedback loops matter",
    "✅ Success requires cross-functional buy-in: finance, compliance, risk, engineering"
])

# Slide 15: Next Steps
add_content_slide(prs, "Your Next Steps", [
    "1. Assess – Profile your top 3 datasets; identify quality gaps",
    "2. Prioritize – Pick one domain (e.g., customers) to pilot",
    "3. Define – Write SLAs, data contracts, and quality rules for pilot",
    "4. Automate – Implement tests and monitoring; set up dashboard",
    "5. Scale – Expand to other domains; formalize governance council and policies"
])

# Save
output_path = r"e:\VScode26\Data_Governance_Quality_Banking.pptx"
prs.save(output_path)
print(f"✅ PowerPoint created: {output_path}")
